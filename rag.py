import streamlit as st
import PyPDF2
import io
import numpy as np
import pickle
import os
import hashlib
from typing import List, Dict, Optional
import re
from sentence_transformers import SentenceTransformer, CrossEncoder
import chromadb
from chromadb.config import Settings
import requests
import json
import docx
from pptx import Presentation
import pandas as pd
from PIL import Image
import pytesseract
from langchain.text_splitter import RecursiveCharacterTextSplitter
import uuid
import time

# Configure Streamlit page
st.set_page_config(
    page_title="Enhanced RAG System with Ollama",
    page_icon="🚀",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Custom CSS for better UI
st.markdown("""
<style>
.chat-message {
    padding: 1rem;
    border-radius: 0.5rem;
    margin-bottom: 1rem;
    display: flex;
    flex-direction: column;
}
.user-message {
    background-color: #2b313e;
    margin-left: 10%;
    border-left: 3px solid #4CAF50;
}
.bot-message {
    background-color: #262730;
    margin-right: 10%;
    border-left: 3px solid #2196F3;
}
.source-box {
    background-color: #f0f2f6;
    padding: 0.5rem;
    border-radius: 0.25rem;
    margin-top: 0.5rem;
    border-left: 3px solid #ff6b6b;
    font-size: 0.9em;
}
.relevance-score {
    background-color: #e3f2fd;
    color: #1976d2;
    padding: 0.2rem 0.5rem;
    border-radius: 1rem;
    font-size: 0.8em;
    font-weight: bold;
}
</style>
""", unsafe_allow_html=True)

class CacheManager:
    """Cache manager for embeddings and responses"""
    
    def __init__(self, cache_dir="./cache"):
        self.cache_dir = cache_dir
        os.makedirs(cache_dir, exist_ok=True)
    
    def get_cache_key(self, text):
        return hashlib.md5(text.encode()).hexdigest()
    
    def cache_embedding(self, text, embedding):
        key = self.get_cache_key(text)
        with open(f"{self.cache_dir}/emb_{key}.pkl", "wb") as f:
            pickle.dump(embedding, f)
    
    def get_cached_embedding(self, text):
        key = self.get_cache_key(text)
        cache_file = f"{self.cache_dir}/emb_{key}.pkl"
        if os.path.exists(cache_file):
            with open(cache_file, "rb") as f:
                return pickle.load(f)
        return None

class AdvancedDocumentProcessor:
    """Enhanced document processor supporting multiple formats"""
    
    def extract_from_pdf(self, pdf_file) -> str:
        try:
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
            return text
        except Exception as e:
            st.error(f"Error reading PDF: {e}")
            return ""
    
    def extract_from_txt(self, txt_file) -> str:
        try:
            content = txt_file.read()
            if isinstance(content, bytes):
                text = content.decode('utf-8')
            else:
                text = content
            return text
        except Exception as e:
            st.error(f"Error reading TXT file: {e}")
            return ""
    
    def extract_from_docx(self, file) -> str:
        try:
            doc = docx.Document(file)
            text = []
            for paragraph in doc.paragraphs:
                text.append(paragraph.text)
            return '\n'.join(text)
        except Exception as e:
            st.error(f"Error reading DOCX file: {e}")
            return ""
    
    def extract_from_pptx(self, file) -> str:
        try:
            prs = Presentation(file)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return '\n'.join(text)
        except Exception as e:
            st.error(f"Error reading PPTX file: {e}")
            return ""
    
    def extract_from_excel(self, file) -> str:
        try:
            df = pd.read_excel(file)
            return df.to_string()
        except Exception as e:
            st.error(f"Error reading Excel file: {e}")
            return ""
    
    def extract_from_image(self, file) -> str:
        try:
            image = Image.open(file)
            text = pytesseract.image_to_string(image)
            return text
        except Exception as e:
            st.error(f"Error extracting text from image: {e}")
            return ""
    
    def process_file(self, file):
        file_type = file.type
        
        if file_type == "application/pdf":
            return self.extract_from_pdf(file)
        elif file_type == "text/plain":
            return self.extract_from_txt(file)
        elif file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            return self.extract_from_docx(file)
        elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            return self.extract_from_pptx(file)
        elif file_type in ["application/vnd.ms-excel", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"]:
            return self.extract_from_excel(file)
        elif file_type.startswith("image/"):
            return self.extract_from_image(file)
        else:
            return None

class AdvancedChunker:
    """Enhanced text chunking with semantic awareness"""
    
    def __init__(self):
        self.recursive_splitter = RecursiveCharacterTextSplitter(
            chunk_size=500,
            chunk_overlap=100,
            separators=["\n\n", "\n", ". ", "! ", "? ", " ", ""]
        )
    
    def chunk_text(self, text, method="recursive"):
        if method == "recursive":
            return self.recursive_splitter.split_text(text)
        else:
            # Fallback to simple chunking
            return self.simple_chunk(text)
    
    def simple_chunk(self, text, chunk_size=500, overlap=100):
        words = text.split()
        chunks = []
        
        for i in range(0, len(words), chunk_size - overlap):
            chunk = ' '.join(words[i:i + chunk_size])
            if chunk.strip() and len(chunk.split()) > 20:
                chunks.append(chunk.strip())
        return chunks

class ChromaVectorStore:
    """Enhanced vector store using ChromaDB with semantic embeddings"""
    
    def __init__(self, collection_name="documents", persist_directory="./chroma_db"):
        self.persist_directory = persist_directory
        os.makedirs(persist_directory, exist_ok=True)
        
        self.client = chromadb.PersistentClient(path=persist_directory)
        self.collection_name = collection_name
        
        # Initialize sentence transformer
        self.encoder = SentenceTransformer('all-MiniLM-L6-v2')
        
        # Get or create collection
        try:
            self.collection = self.client.get_collection(name=collection_name)
        except:
            self.collection = self.client.create_collection(name=collection_name)
        
        self.cache_manager = CacheManager()
    
    def encode_text(self, text):
        """Encode text with caching"""
        cached = self.cache_manager.get_cached_embedding(text)
        if cached is not None:
            return cached
        
        embedding = self.encoder.encode(text)
        self.cache_manager.cache_embedding(text, embedding)
        return embedding
    
    def add_documents(self, chunks, metadata):
        """Add documents to the vector store"""
        ids = [str(uuid.uuid4()) for _ in chunks]
        
        # Get embeddings
        embeddings = [self.encode_text(chunk).tolist() for chunk in chunks]
        
        self.collection.add(
            embeddings=embeddings,
            documents=chunks,
            metadatas=metadata,
            ids=ids
        )
        
        return ids
    
    def search(self, query, k=10):
        """Search for similar documents"""
        query_embedding = self.encode_text(query).tolist()
        
        results = self.collection.query(
            query_embeddings=[query_embedding],
            n_results=k,
            include=['documents', 'metadatas', 'distances']
        )
        
        search_results = []
        if results['documents'] and results['documents'][0]:
            for i, (doc, metadata, distance) in enumerate(zip(
                results['documents'][0],
                results['metadatas'][0],
                results['distances'][0]
            )):
                # Convert distance to similarity score
                similarity = 1 / (1 + distance)
                search_results.append({
                    'chunk': doc,
                    'score': similarity,
                    'metadata': metadata
                })
        
        return search_results
    
    def clear_collection(self):
        """Clear all documents from the collection"""
        try:
            self.client.delete_collection(name=self.collection_name)
            self.collection = self.client.create_collection(name=self.collection_name)
        except:
            pass

class Reranker:
    """Cross-encoder reranking for improved relevance"""
    
    def __init__(self, model_name="cross-encoder/ms-marco-MiniLM-L-6-v2"):
        try:
            self.reranker = CrossEncoder(model_name)
            self.available = True
        except:
            self.available = False
            st.warning("Reranker model not available. Install sentence-transformers with: pip install sentence-transformers")
    
    def rerank(self, query, search_results, top_k=5):
        if not self.available or not search_results:
            return search_results[:top_k]
        
        try:
            # Prepare query-document pairs
            pairs = [(query, result['chunk']) for result in search_results]
            
            # Get reranking scores
            scores = self.reranker.predict(pairs)
            
            # Combine with original results
            for i, result in enumerate(search_results):
                result['rerank_score'] = float(scores[i])
            
            # Sort by reranking score
            reranked = sorted(search_results, key=lambda x: x['rerank_score'], reverse=True)
            return reranked[:top_k]
        except:
            return search_results[:top_k]

class QueryEnhancer:
    """Enhanced query processing and expansion"""
    
    def __init__(self, ollama_provider):
        self.ollama = ollama_provider
    
    def expand_query(self, query):
        """Generate related questions/keywords"""
        if not self.ollama or not self.ollama.connected:
            return [query]
        
        prompt = f"""Generate 2 related questions or alternative phrasings for this query. Keep them concise and relevant:
        
Original: {query}

Alternative 1:"""
        
        try:
            expansion = self.ollama.generate_response(prompt, max_tokens=100)
            alternatives = [line.strip() for line in expansion.split('\n') if line.strip() and len(line.strip()) > 10]
            return [query] + alternatives[:2]
        except:
            return [query]
    
    def extract_query_parts(self, query):
        """Extract multiple questions from complex queries"""
        parts = []
        
        # Split by question marks
        q_parts = re.split(r'\?', query)
        for part in q_parts:
            part = part.strip()
            if part:
                # Further split by conjunctions
                sub_parts = re.split(r'\b(?:and|also|additionally|furthermore|moreover)\b', part, flags=re.IGNORECASE)
                for sub_part in sub_parts:
                    sub_part = sub_part.strip()
                    if len(sub_part) > 10:
                        parts.append(sub_part)
        
        return parts if parts else [query]

class OllamaProvider:
    """Enhanced Ollama LLM provider"""
    
    def __init__(self, model: str = "llama2", base_url: str = "http://localhost:11434"):
        self.name = "Ollama"
        self.model = model
        self.base_url = base_url
        self.generate_url = f"{base_url}/api/generate"
        self.tags_url = f"{base_url}/api/tags"
        self.connected = False
        
    def test_connection(self) -> bool:
        try:
            response = requests.get(self.tags_url, timeout=5)
            if response.status_code == 200:
                self.connected = True
                return True
            else:
                self.connected = False
                return False
        except Exception:
            self.connected = False
            return False
    
    def get_available_models(self) -> List[str]:
        try:
            if not self.test_connection():
                return []
            
            response = requests.get(self.tags_url, timeout=5)
            if response.status_code == 200:
                data = response.json()
                models = [model['name'] for model in data.get('models', [])]
                return models
            return []
        except Exception:
            return []
    
    def generate_response(self, prompt: str, max_tokens: int = 1000) -> str:
        try:
            data = {
                "model": self.model,
                "prompt": prompt,
                "stream": False,
                "options": {
                    "num_predict": max_tokens,
                    "temperature": 0.7,
                    "top_p": 0.9,
                }
            }
            
            response = requests.post(self.generate_url, json=data, timeout=120)
            response.raise_for_status()
            
            result = response.json()
            return result.get("response", "").strip()
            
        except requests.exceptions.Timeout:
            return "Error: Request timed out."
        except requests.exceptions.ConnectionError:
            return "Error: Cannot connect to Ollama."
        except Exception as e:
            return f"Error: {str(e)}"

class EnhancedRAG:
    """Enhanced RAG system with all improvements"""
    
    def __init__(self):
        self.vector_store = ChromaVectorStore()
        self.document_processor = AdvancedDocumentProcessor()
        self.chunker = AdvancedChunker()
        self.reranker = Reranker()
        self.ollama = None
        self.query_enhancer = None
        self.documents = []
        
    def setup_ollama(self, model: str, base_url: str = "http://localhost:11434"):
        """Setup Ollama provider and query enhancer"""
        self.ollama = OllamaProvider(model, base_url)
        connection_status = self.ollama.test_connection()
        if connection_status:
            self.query_enhancer = QueryEnhancer(self.ollama)
        return connection_status
    
    def add_document(self, filename: str, file_obj):
        """Add a document to the RAG system"""
        # Extract text based on file type
        text = self.document_processor.process_file(file_obj)
        
        if not text or not text.strip():
            st.warning(f"Could not extract text from {filename}!")
            return False
        
        # Clean and chunk text
        cleaned_text = re.sub(r'\s+', ' ', text.strip())
        chunks = self.chunker.chunk_text(cleaned_text)
        
        if not chunks:
            st.warning("No text chunks created from document!")
            return False
        
        # Create metadata for each chunk
        metadata = []
        for i, chunk in enumerate(chunks):
            metadata.append({
                'filename': filename,
                'chunk_id': i,
                'doc_id': len(self.documents),
                'chunk_length': len(chunk)
            })
        
        # Add to vector store
        try:
            ids = self.vector_store.add_documents(chunks, metadata)
            
            # Store document info
            doc_info = {
                'filename': filename,
                'text': cleaned_text,
                'chunk_count': len(chunks),
                'chunk_ids': ids
            }
            self.documents.append(doc_info)
            
            st.success(f"✅ Added '{filename}' with {len(chunks)} chunks to the knowledge base!")
            return True
            
        except Exception as e:
            st.error(f"Error adding document: {e}")
            return False
    
    def search_and_rerank(self, query: str, top_k: int = 8) -> List[Dict]:
        """Enhanced search with query expansion and reranking"""
        all_results = []
        
        # Use query enhancer if available
        if self.query_enhancer:
            queries = self.query_enhancer.expand_query(query)
        else:
            queries = [query]
        
        # Search with multiple query variants
        for q in queries[:3]:  # Limit to 3 query variants
            results = self.vector_store.search(q, k=top_k)
            all_results.extend(results)
        
        # Remove duplicates based on chunk content
        unique_results = {}
        for result in all_results:
            chunk_key = result['chunk'][:100]  # Use first 100 chars as key
            if chunk_key not in unique_results or result['score'] > unique_results[chunk_key]['score']:
                unique_results[chunk_key] = result
        
        final_results = list(unique_results.values())
        
        # Rerank if available
        if self.reranker.available and final_results:
            final_results = self.reranker.rerank(query, final_results, top_k)
        
        return sorted(final_results, key=lambda x: x.get('rerank_score', x['score']), reverse=True)[:top_k]
    
    def create_enhanced_prompt(self, query: str, context_chunks: List[str], chat_history: List[Dict] = None) -> str:
        """Create enhanced prompt with better context organization"""
        context = "\n\n---DOCUMENT SECTION---\n\n".join(context_chunks[:5])
        
        conversation = ""
        if chat_history:
            for msg in chat_history[-6:]:
                if msg['role'] == 'user':
                    conversation += f"Human: {msg['content']}\n"
                else:
                    conversation += f"Assistant: {msg['content']}\n"
        
        prompt = f"""You are an expert AI assistant that provides accurate, detailed answers based on document context. You excel at synthesizing information from multiple sources.

DOCUMENT CONTEXT:
{context}

CONVERSATION HISTORY:
{conversation}

CURRENT QUESTION: {query}

INSTRUCTIONS:
- Provide a comprehensive answer using ONLY the information from the document context
- Synthesize information from multiple document sections when relevant
- Include specific details and examples from the documents
- If information is incomplete, state what you can answer and what's missing
- Use clear structure with bullet points or numbered lists when appropriate
- Cite specific sections when making claims

ANSWER:"""
        return prompt
    
    def generate_answer_with_sources(self, query: str, chat_history: List[Dict] = None) -> Dict:
        """Generate answer with source citations and relevance scores"""
        # Search for relevant chunks
        search_results = self.search_and_rerank(query, top_k=8)
        
        if not search_results:
            return {
                'answer': "I couldn't find relevant information in the uploaded documents to answer your question.",
                'sources': [],
                'search_results': []
            }
        
        # Extract chunks for context
        context_chunks = [result['chunk'] for result in search_results[:5]]
        
        # Generate answer
        if self.ollama and self.ollama.connected:
            try:
                prompt = self.create_enhanced_prompt(query, context_chunks, chat_history)
                answer = self.ollama.generate_response(prompt, max_tokens=1500)
                
                if answer.startswith("Error:"):
                    answer = self.create_fallback_answer(context_chunks)
                    
            except Exception as e:
                answer = self.create_fallback_answer(context_chunks)
        else:
            answer = self.create_fallback_answer(context_chunks)
        
        # Prepare sources
        sources = []
        for i, result in enumerate(search_results[:5]):
            sources.append({
                'id': i + 1,
                'filename': result['metadata']['filename'],
                'score': result['score'],
                'rerank_score': result.get('rerank_score', 0),
                'preview': result['chunk'][:150] + "...",
                'chunk_id': result['metadata'].get('chunk_id', 0)
            })
        
        return {
            'answer': answer,
            'sources': sources,
            'search_results': search_results
        }
    
    def create_fallback_answer(self, context_chunks: List[str]) -> str:
        """Create fallback answer when Ollama is not available"""
        context = "\n\n".join(context_chunks[:3])
        return f"**Based on the uploaded documents:**\n\n{context[:1500]}{'...' if len(context) > 1500 else ''}\n\n*Note: This is a direct excerpt. Ollama connection not available for enhanced processing.*"
    
    def clear_all_documents(self):
        """Clear all documents from the system"""
        self.vector_store.clear_collection()
        self.documents = []

def setup_ollama_config():
    """Enhanced Ollama configuration UI"""
    st.subheader("🦙 Ollama Configuration")
    
    base_url = st.text_input(
        "Ollama Base URL:",
        value="http://localhost:11434",
        help="URL where Ollama is running"
    )
    
    test_ollama = OllamaProvider(base_url=base_url)
    
    col1, col2 = st.columns([1, 1])
    
    with col1:
        if st.button("🔌 Test Connection", type="secondary"):
            with st.spinner("Testing connection..."):
                if test_ollama.test_connection():
                    st.success("✅ Connected!")
                    models = test_ollama.get_available_models()
                    if models:
                        st.write("**Available models:**")
                        for model in models:
                            st.write(f"• {model}")
                else:
                    st.error("❌ Connection failed")
    
    with col2:
        st.caption("**Quick Setup:**")
        st.code("ollama pull llama2\nollama serve", language="bash")
    
    # Model selection
    available_models = test_ollama.get_available_models()
    
    if available_models:
        model = st.selectbox("Select Model:", available_models)
    else:
        model = st.text_input("Model Name:", value="llama2")
    
    # Connection status indicator
    if test_ollama.test_connection():
        st.success("🟢 Ollama Online")
        return model, base_url, True
    else:
        st.error("🔴 Ollama Offline")
        return model, base_url, False

def display_enhanced_chat():
    """Enhanced chat interface with sources"""
    st.header("💬 Chat with Your Documents")
    
    # Chat container
    chat_container = st.container()
    
    with chat_container:
        for message in st.session_state.chat_history:
            if message['role'] == 'user':
                st.markdown(f"""
                <div class="chat-message user-message">
                    <strong>You:</strong> {message['content']}
                </div>
                """, unsafe_allow_html=True)
            else:
                st.markdown(f"""
                <div class="chat-message bot-message">
                    <strong>🚀 Assistant:</strong> {message['content']}
                </div>
                """, unsafe_allow_html=True)
                
                # Show sources if available
                if 'sources' in message and message['sources']:
                    with st.expander("📚 Sources Used", expanded=False):
                        for source in message['sources']:
                            relevance_color = "🟢" if source['score'] > 0.7 else "🟡" if source['score'] > 0.4 else "🔴"
                            st.markdown(f"""
                            <div class="source-box">
                                <strong>{relevance_color} Source {source['id']}: {source['filename']}</strong><br>
                                <span class="relevance-score">Relevance: {source['score']:.3f}</span>
                                {f" | Rerank: {source['rerank_score']:.3f}" if source.get('rerank_score', 0) > 0 else ""}<br>
                                <em>{source['preview']}</em>
                            </div>
                            """, unsafe_allow_html=True)

def main():
    # Header
    st.title("🚀 Enhanced RAG System with Ollama")
    st.markdown("**Advanced AI-powered document chat with semantic search, reranking, and source citations**")
    
    # Initialize enhanced RAG system
    if 'rag_system' not in st.session_state:
        with st.spinner("Initializing enhanced RAG system..."):
            st.session_state.rag_system = EnhancedRAG()
    
    # Initialize chat history
    if 'chat_history' not in st.session_state:
        st.session_state.chat_history = []
    
    if 'last_input' not in st.session_state:
        st.session_state.last_input = ""
    
    rag = st.session_state.rag_system
    
    # Sidebar configuration
    with st.sidebar:
        st.header("⚙️ Configuration")
        
        # Ollama setup
        with st.expander("🦙 Ollama Settings", expanded=True):
            model, base_url, is_connected = setup_ollama_config()
            if is_connected:
                rag.setup_ollama(model, base_url)
        
        st.divider()
        
        # Document management
        st.header("📁 Document Management")
        
        uploaded_files = st.file_uploader(
            "Upload Documents",
            type=['pdf', 'txt', 'docx', 'pptx', 'xlsx', 'xls', 'png', 'jpg', 'jpeg'],
            accept_multiple_files=True,
            help="Supports: PDF, TXT, DOCX, PPTX, Excel, Images"
        )
        
        if uploaded_files:
            for uploaded_file in uploaded_files:
                existing_docs = [doc['filename'] for doc in rag.documents]
                if uploaded_file.name not in existing_docs:
                    with st.spinner(f"Processing {uploaded_file.name}..."):
                        rag.add_document(uploaded_file.name, uploaded_file)
        
        # Document display
        if rag.documents:
            st.subheader(f"📋 Loaded Documents ({len(rag.documents)})")
            for i, doc in enumerate(rag.documents):
                with st.expander(f"📄 {doc['filename']}", expanded=False):
                    st.metric("Chunks", doc['chunk_count'])
                    st.text_area(
                        "Preview", 
                        doc['text'][:300] + "...", 
                        height=100, 
                        key=f"preview_{i}",
                        disabled=True
                    )
        else:
            st.info("Upload documents to start chatting!")
        
        # Clear button
        if st.button("🗑️ Clear All Documents", type="secondary"):
            rag.clear_all_documents()
            st.session_state.chat_history = []
            st.session_state.last_input = ""
            st.rerun()
    
    # Main chat interface
    display_enhanced_chat()
    
    # Query input
    with st.form(key="enhanced_chat_form", clear_on_submit=True):
        col1, col2 = st.columns([4, 1])
        
        with col1:
            user_query = st.text_area(
                "Ask about your documents:",
                placeholder="e.g., What are the main findings? Also, what methodology was used?",
                height=100,
                key="user_input_enhanced"
            )
        
        with col2:
            st.write("")  # Spacing
            st.write("")  # Spacing
            submit_button = st.form_submit_button("🚀 Send", type="primary", use_container_width=True)
            
            # Query enhancement options
            enhance_query = st.checkbox("🔍 Enhance query", value=True, help="Use query expansion for better results")
    
    # Process query
    if submit_button and user_query and user_query != st.session_state.last_input:
        if rag.documents:
            st.session_state.last_input = user_query
            
            # Add user message
            st.session_state.chat_history.append({
                'role': 'user',
                'content': user_query
            })
            
            with st.spinner("🔍 Searching documents and generating comprehensive answer..."):
                start_time = time.time()
                
                # Generate enhanced answer
                result = rag.generate_answer_with_sources(
                    user_query, 
                    st.session_state.chat_history[:-1]
                )
                
                processing_time = time.time() - start_time
                
                # Add assistant message with sources
                assistant_message = {
                    'role': 'assistant',
                    'content': result['answer'],
                    'sources': result['sources'],
                    'processing_time': processing_time
                }
                
                st.session_state.chat_history.append(assistant_message)
            
            st.rerun()
        
        elif user_query:
            st.warning("Please upload some documents first!")
    
    # System status
    st.sidebar.divider()
    st.sidebar.subheader("📊 System Status")
    
    if rag.ollama and rag.ollama.connected:
        st.sidebar.success(f"🦙 Ollama: {rag.ollama.model}")
    else:
        st.sidebar.error("🦙 Ollama: Disconnected")
    
    if rag.reranker.available:
        st.sidebar.success("🔄 Reranker: Active")
    else:
        st.sidebar.warning("🔄 Reranker: Unavailable")
    
    st.sidebar.info(f"📚 Documents: {len(rag.documents)}")
    
    # Performance tips
    with st.sidebar.expander("💡 Tips", expanded=False):
        st.markdown("""
        **For best results:**
        - Upload related documents
        - Ask specific questions
        - Use natural language
        - Check source citations
        - Try multiple query phrasings
        """)

if __name__ == "__main__":
    main()